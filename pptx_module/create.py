from pptx import Presentation 

class PPTCreator:
    def __init__(self, path):
        self.path = path
        self.prs = Presentation()
        self.data= [
            "\nSlide 1\nTemp_pptx created.\nNext moving to slide 2.",
            "\nSlide 2\nJumping to slide 4.",
            "\nSlide 3\nJumping to end slide.",
            "\nSlide 4\nGoing backward",
            "\nSlide 5\nClosing the presantation"
        ]

    def add_slide(self, slide_layout):
        slide = self.prs.slides.add_slide(slide_layout)
        return slide

    def add_text(self, slide, text):
        placeholder = slide.shapes.title or slide.placeholders[1]
        placeholder.text = text


    def main(self):
        
        print("Adding slides and text in ppt.")
        slide_layout = self.prs.slide_layouts[5]

        for each in self.data:
            self.add_text(self.add_slide(slide_layout), each)
            
        self.prs.save(self.path)
        print(f"Presentation '{self.path}' created successfully.")

if __name__ == "__main__":
    creator = PPTCreator("test.pptx")
    creator.main()
