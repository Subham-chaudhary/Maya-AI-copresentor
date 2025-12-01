import json
import os
from pptx import Presentation

class PPTParser:
    def __init__(self, file_path):
        self.file_name = os.path.splitext(os.path.basename(file_path))[0]
        self.prs = Presentation(file_path)


    def toPX(self, num):
        return round(num/1000)

    def parse(self):
        slides_data = [ {"file_name": self.file_name}]

        for idx, slide in enumerate(self.prs.slides, start=1):
            slide_info = {
                "slide_number": idx,
                "title": slide.shapes.title.text if slide.shapes.title else None,
                "shapes": []
            }

            for shape in slide.shapes:
                shape_info = {
                    "type": shape.__class__.__name__,
                    "left": self.toPX(shape.left), 
                    "top": self.toPX(shape.top),
                    "width": self.toPX(shape.width),
                    "height": self.toPX(shape.height),
                }

                if shape.has_text_frame:
                    text = "\n".join(
                        para.text for para in shape.text_frame.paragraphs
                    )
                    shape_info["text"] = text
                else:
                    shape_info["text"] = None

                slide_info["shapes"].append(shape_info)

            slides_data.append(slide_info)

        return slides_data


if __name__ == "__main__":
    file_path = "test."
    parser = PPTParser(file_path)
    data = parser.parse()
    with open(os.path.splitext(os.path.basename(file_path))[0]+".json", "w") as outfile:
        json.dump(data, outfile, indent=4)
